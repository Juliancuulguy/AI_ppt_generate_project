from pathlib import Path
from tkinter import Tk, Canvas, Entry, Text, Button, PhotoImage
from tkinter import *
from tkinter.ttk import *
from tkinter import messagebox, ttk
from tkinter.colorchooser import askcolor
import tkinter.font as font
from tkinter.constants import *
from tkinter.constants import CENTER
from dotenv import dotenv_values
import openai
import customtkinter as ctk
import tkinter as tk
import re
import openai,collections,_collections_abc,json,io, os, numpy, getpass, warnings, random, sys
from pptx import Presentation,util
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches,Pt,Cm
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE,PP_ALIGN
from zhon.hanzi import punctuation
from dotenv import dotenv_values
from IPython.display import display
from PIL import Image
from stability_sdk import client
import stability_sdk.interfaces.gooseai.generation.generation_pb2 as generation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.dml.fill import (
    CT_BlipFillProperties,
    CT_GradientFillProperties,
    CT_GroupFillProperties,
    CT_NoFillProperties,
    CT_PatternFillProperties,
    CT_SolidColorFillProperties,
)
#import aspose.slides as slides
#import aspose.pydrawing as drawing
import random
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from PyQt5.QtWidgets import *

OUTPUT_PATH = Path(__file__).parent

ansl = []
pic_loc = []
con_word = []

image_num = 0
max_image = 3
# 標題編號 page2說明 按鈕說明 pagethree 說明
#-----keys

openai.api_key=''

os.environ['STABILITY_HOST'] = 'grpc.stability.ai:443'
os.environ['STABILITY_KEY'] = ''

#-----PPT page
page_of_pic = ["ppt1.png", "ppt2.png", "ppt3.png", "ppt4.png",
               "ppt5.png", "ppt6.png", "ppt7.png", "ppt8.png",
               "ppt9.png", "ppt10.png", "ppt11.png", "ppt12.png"]
pic_path = ["self.image_0", "self.image_1", "self.image_2", "self.image_3"]
pic_pic = ["self.img0", "self.img1", "self.img2", "self.img3",
           "self.img4", "self.img5","self.img6", "self.img7",
           "self.img8", "self.img9","self.img10", "self.img11"]
layout = 0
picpath = 0

ppt_page_index = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12]
pindex = 0

pagechangei = 0
amounts = 0

pptchoose = 0 # 選擇的模板

cont = 0 # 模板中的內文數
picc = 0 # 模板中的圖片數

img_list = [] # 選取的圖片放置處
imgchoose = 0 # 選取圖片數

#-----chatGPT's list
key_word = ""
chatans = []            #chatGPT's answer
spl = [] #生成的標題(第0、1個是空格，從第2個開始才是標題)
spl2 = []
spl_copy = []
spp3 = []
sub_word = []
try1 = [] 
contents_done = [] #生成的內文            
spp2 = ''
sub = ''
tousr = ''
amounts = 0 #生成的標題數量
page_of_ppt = 0
tindex = 0
consave = []
hint_yn = 0
note = '' #簡報備忘錄
regenerate = False #重新生成 
do = 0 #do_QA的temperature

mc1, mc2, mc3 = 112,94,81 #主頁面模板色
mcname = ''
mtc1, mtc2, mtc3 = 80,76,68 #主頁面文字色
mtcname = ''
sc1, sc2, sc3 = 112,94,81 #內文頁模板色
scname = ''
stc1, stc2, stc3 = 80,76,68 #內文頁文字色
stcname = ''

prs = Presentation()

def next():
    global page_of_ppt
    page_of_ppt += 1

def change_regenerate():
    global regenerate
    regenerate = False
#--chatGPT - prompt dealing
def doQA(prompts, max_tokens, temperature):
    text = openai.Completion.create(
            model="text-davinci-003",
            prompt=prompts,
            max_tokens=max_tokens, 
            temperature=temperature
    )
    chatans.append(text['choices'][0]['text'])

def temperature_up():
    global do
    do = do + 0.1
    if do == 0.9:
        temperature_0()

def temperature_0():
    global do
    do = 0

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))
    
def _set_line_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))

def generate_ppt():
    global try1, spl, img_list, key_word, prs, amounts, mc1, mc2, mc3, mtc1, mtc2, mtc3    # 內文, 標題, 圖片
    #開啟新的簡報物件

    c_1_r, c_1_g, c_1_b = 253, 253, 253

    #第一頁模板(封面)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
    #背景
    slide_bkg = slide.background
    slide_bkg_fill =  slide_bkg.fill
    slide_bkg_fill.solid()
    slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 

    #add rectangle
    left = util.Cm(0)
    top = util.Cm(0)
    width = util.Cm(4)
    height = util.Cm(19.05)
    #shapes = slide.shapes
    rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rectangle_1.fill.solid()
    rectangle_1.fill.fore_color.rgb = RGBColor(mc1, mc2, mc3)
    #rectangle.line
    rectangle_1.line.color.rgb = RGBColor(mc1, mc2, mc3)

    #標題
    left = util.Cm(5.5)
    top = util.Cm(7.5)
    width = util.Cm(18)
    height = util.Cm(3)
    #shapes = slide.shapes
    RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    RECTANGLE_1.fill.solid()
    RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

    tf = RECTANGLE_1.text_frame
    p = tf.paragraphs[0]
    p.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)
    p.font.size = Pt(32) #大小
    p.font.name = 'Arial' #字體
    p.font.bold = True #加粗 
    p.alignment = PP_ALIGN.LEFT
    p.text = f'{key_word}'

    RECTANGLE_1.line.fill.background()

    _set_shape_transparency(RECTANGLE_1,0)

    #add line
    start_x = util.Cm(5.95) 
    start_y = util.Cm(10)
    end_x = util.Cm(14.95)
    end_y = util.Cm(10)

    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
    line1.line.fill.background()
    line1.line.fill.solid()
    line1.line.fill.fore_color.rgb = RGBColor(0, 0, 0) #黑色

    #副標題
    left = util.Cm(5.5) 
    top = util.Cm(10.5)
    width = util.Cm(18)
    height = util.Cm(3)
    #shapes = slide.shapes
    RECTANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    RECTANGLE_2.fill.solid()
    RECTANGLE_2.fill.fore_color.rgb = RGBColor(253,253,253)
    
    chatans.clear()
    prompts = f'我要製作一個簡報，主題是關於{key_word}，請以一簡短句生成副標題'
    doQA(prompts, 512, 0.3)
    sub = ''.join(chatans[0])
    sub_word = sub.split('\n')

    tf = RECTANGLE_2.text_frame
    p = tf.paragraphs[0]
    p.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)
    p.font.size = Pt(28) #大小
    p.font.name = 'Arial' #字體
    p.font.bold = False #加粗 
    p.alignment = PP_ALIGN.LEFT
    p.text = f'{sub_word[2]}'

    RECTANGLE_2.line.fill.background()

    _set_shape_transparency(RECTANGLE_2,0)

    #第二頁模板(目錄頁,3個標題) 
    slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
    #背警
    slide_bkg = slide.background
    slide_bkg_fill =  slide_bkg.fill
    slide_bkg_fill.solid()
    slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 
    i = 0 
    j = 0
    k = 1

    #目錄
    left = util.Cm(11.35)
    top = util.Cm(1.55)
    width = util.Cm(3)
    height = util.Cm(3)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)# 灰色
    p.font.size = Pt(32) #大小
    p.font.name = 'Arial' #字體
    p.font.bold = True #加粗
    p.text='目錄'

    if len(spl) > 3:
        for i in range(len(spl)):
            left = util.Cm(3.85) 
            top = util.Cm(4+i*2)
            width = util.Cm(3)
            height = util.Cm(3)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            tf = text_box.text_frame
            p2 = tf.paragraphs[0]
            p2.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)# 灰色
            p2.font.size = Pt(28) #大小
            p2.font.name = 'Arial' #字體
            p2.font.bold = True #加粗
            p2 = tf.paragraphs[0]
            p2.text = f"{j}"+f"{k}"
            
            k+=1

        #標題1
            left = util.Cm(6) 
            top = util.Cm(3.2 + i*2)
            width = util.Cm(15)
            height = util.Cm(3)
            #shapes = slide.shapes
            RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            RECTANGLE1.fill.solid()
            RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

            tf = RECTANGLE1.text_frame
            p1 = tf.paragraphs[0]
            p1.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)
            p1.font.size = Pt(18) #大小
            p1.font.name = 'Arial' #字體
            p1.font.bold = False #加粗 
            p1.alignment = PP_ALIGN.LEFT
            p1.text = f'{spl[i]}'

            RECTANGLE1.line.fill.background()

            _set_shape_transparency(RECTANGLE1,0)
    else:
        for i in range(len(spl)):
            
            left = util.Cm(3.85) 
            top = util.Cm(5+i*4)
            width = util.Cm(3)
            height = util.Cm(3)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            tf = text_box.text_frame
            p2 = tf.paragraphs[0]
            p2.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)# 灰色
            p2.font.size = Pt(32) #大小
            p2.font.name = 'Arial' #字體
            p2.font.bold = True #加粗
            p2 = tf.paragraphs[0]
            p2.text = f"{j}"+f"{k}"
            
            k+=1

        #標題1
            left = util.Cm(6) 
            top = util.Cm(4.2 + i*4)
            width = util.Cm(15)
            height = util.Cm(3)
            #shapes = slide.shapes
            RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            RECTANGLE1.fill.solid()
            RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

            tf = RECTANGLE1.text_frame
            p1 = tf.paragraphs[0]
            p1.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)
            p1.font.size = Pt(20) #大小
            p1.font.name = 'Arial' #字體
            p1.font.bold = False #加粗 
            p1.alignment = PP_ALIGN.LEFT
            p1.text = f'{spl[i]}'

            RECTANGLE1.line.fill.background()

            _set_shape_transparency(RECTANGLE1,0)

def generate_one_page_ppt():
    global prs, try1, spl_copy, img_list, pptchoose, contents_done, sc1, sc2, sc3, stc1, stc2, stc3, note    # PPT檔案, 內文, 標題, 圖片, 選擇的模板
    #開啟新的簡報物件
    c_1_r, c_1_g, c_1_b = 253, 253, 253

    if pptchoose == '1':
        #第三頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 
        #add rectangle
        left = util.Cm(0)
        top = util.Cm(0)
        width = util.Cm(15.55)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add ROUNDED_RECTANGLE
        left = util.Cm(12.55)
        top = util.Cm(3.55)
        width = util.Cm(11)
        height = util.Cm(11.5)
        #shapes = slide.shapes
        rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rounded_rectangle.fill.solid()
        rounded_rectangle.fill.fore_color.rgb = RGBColor(241,237,233)
        #rectangle.line
        rounded_rectangle.line.color.rgb = RGBColor(241,237,233)

        #貼圖片
        tf = rounded_rectangle.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(14) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        #add title
        left = util.Cm(1.5) 
        top = util.Cm(2)
        width = util.Cm(10)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE.fill.solid()
        RECTANGLE.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)

        tf = RECTANGLE.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE.line.fill.background()

        _set_shape_transparency(RECTANGLE,0)


        #add content
        left = util.Cm(1.5) 
        top = util.Cm(6)
        width = util.Cm(10)
        height = util.Cm(9)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[0]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '2':
        #第四頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 
        #add rectangle
        left = util.Cm(0)
        top = util.Cm(13.5)
        width = util.Cm(25.38)
        height = util.Cm(5.55)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add ROUNDED_RECTANGLE
        left = util.Cm(12.55)
        top = util.Cm(3.55)
        width = util.Cm(11)
        height = util.Cm(11.5)
        #shapes = slide.shapes
        rounded_rectangle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rounded_rectangle.fill.solid()
        rounded_rectangle.fill.fore_color.rgb = RGBColor(241,237,233)
        #rectangle.line
        rounded_rectangle.line.color.rgb = RGBColor(241,237,233)

        #貼圖片
        tf = rounded_rectangle.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        #add title
        left = util.Cm(1.5) 
        top = util.Cm(2)
        width = util.Cm(10)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE.fill.solid()
        RECTANGLE.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)

        tf = RECTANGLE.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE.line.fill.background()

        _set_shape_transparency(RECTANGLE,0)


        #add content
        left = util.Cm(1.5) 
        top = util.Cm(6)
        width = util.Cm(10)
        height = util.Cm(6.5)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[0]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note


    if pptchoose == '3':
        #第五頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 

        #add title
        left = util.Cm(13) 
        top = util.Cm(1)
        width = util.Cm(10)
        height = util.Cm(2)
        #shapes = slide.shapes
        RECTANGLE = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE.fill.solid()
        RECTANGLE.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE.line.fill.background()

        _set_shape_transparency(RECTANGLE,0)


        #add rectangle
        left = util.Cm(3.5)
        top = util.Cm(0)
        width = util.Cm(5)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        # add ROUNDED_RECTANGLE1
        left = util.Cm(5.5) 
        top = util.Cm(2.75)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        rounded_rectangle1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rounded_rectangle1.fill.solid()
        rounded_rectangle1.fill.fore_color.rgb = RGBColor(241,237,237)
        #rectangle.line
        rounded_rectangle1.line.color.rgb = RGBColor(241,237,237)
        #內文
        tf = rounded_rectangle1.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        # add ROUNDED_RECTANGLE2
        left = util.Cm(5.5) 
        top = util.Cm(10.75)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        rounded_rectangle2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rounded_rectangle2.fill.solid()
        rounded_rectangle2.fill.fore_color.rgb = RGBColor(241,237,237)
        #rectangle.line
        rounded_rectangle2.line.color.rgb = RGBColor(241,237,237)
        #內文
        tf = rounded_rectangle2.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[1]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        #add oval
        i1 = 0 
        for i1 in range(2):

            left = util.Cm(13) 
            top = util.Cm(5.5+i1*7.5)
            width = util.Cm(1)
            height = util.Cm(1)
        #shapes = slide.shapes
            Oval_1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            Oval_1.fill.solid()
            Oval_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
            Oval_1.line.color.rgb = RGBColor(sc1, sc2, sc3)
            

        #add content1
        left = util.Cm(15) 
        top = util.Cm(3.5)
        width = util.Cm(8)
        height = util.Cm(6)
        #shapes = slide.shapes
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add content2
        left = util.Cm(15) 
        top = util.Cm(11)
        width = util.Cm(8)
        height = util.Cm(6)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note


    if pptchoose == '4':
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 

        #add rectangle
        left = util.Cm(0)
        top = util.Cm(0)
        width = util.Cm(4)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add title
        left = util.Cm(5.5) 
        top = util.Cm(1.25)
        width = util.Cm(18)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE.fill.solid()
        RECTANGLE.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE.line.fill.background()

        _set_shape_transparency(RECTANGLE,0)


        #add line
        start_x = util.Cm(5.95) 
        start_y = util.Cm(3.8)
        end_x = util.Cm(14.95)
        end_y = util.Cm(3.8)

        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
        line1.line.fill.background()
        line1.line.fill.solid()
        line1.line.fill.fore_color.rgb = RGBColor(0, 0, 0) #黑色

        #add content1
        left = util.Cm(5.5) 
        top = util.Cm(5.5)
        width = util.Cm(14)
        height = util.Cm(5)
        #shapes = slide.shapes
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.line_spacing = 1.0 
        p1.alignment = PP_ALIGN.LEFT

        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add content2
        left = util.Cm(5.5) 
        top = util.Cm(11.5)
        width = util.Cm(14)
        height = util.Cm(5)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '5':
        #第七頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 

        #add rectangle
        left = util.Cm(0)
        top = util.Cm(0)
        width = util.Cm(4)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add title
        left = util.Cm(5.5) 
        top = util.Cm(1.25)
        width = util.Cm(18)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE.fill.solid()
        RECTANGLE.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE.line.fill.background()

        _set_shape_transparency(RECTANGLE,0)

        #add line
        start_x = util.Cm(5.95) 
        start_y = util.Cm(3.8)
        end_x = util.Cm(14.95)
        end_y = util.Cm(3.8)

        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
        line1.line.fill.background()
        line1.line.fill.solid()
        line1.line.fill.fore_color.rgb = RGBColor(0, 0, 0) #黑色

        #add content1
        left = util.Cm(5.5) 
        top = util.Cm(5)
        width = util.Cm(15)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add content2
        left = util.Cm(5.5) 
        top = util.Cm(9)
        width = util.Cm(15)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        #add content3
        left = util.Cm(5.5) 
        top = util.Cm(13)
        width = util.Cm(15)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE3.fill.solid()
        RECTANGLE3.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE3.text_frame
        p3 = tf.paragraphs[0]
        p3.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p3.font.size = Pt(20) #大小
        p3.font.name = 'Arial' #字體
        p3.font.bold = False #加粗 
        p3.alignment = PP_ALIGN.LEFT
        p3.line_spacing = 1.0 
        p3.text = f'{contents_done[2]}'

        RECTANGLE3.line.fill.background()

        _set_shape_transparency(RECTANGLE3,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '6':
        #第八頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b ) 
        #add rectangle
        left = util.Cm(3.2)
        top = util.Cm(-3.19)
        width = util.Cm(18.9)
        height = util.Cm(25.6)
        #shapes = slide.shapes
        RIGHT_TRIANGLE = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE.fill.solid()
        RIGHT_TRIANGLE.fill.fore_color.rgb = RGBColor(224,224,224)
        #RIGHT_TRIANGLE_1.line
        #RIGHT_TRIANGLE.line.color.rgb = RGBColor(224,224,244)

        RIGHT_TRIANGLE.line.fill.background()

        RIGHT_TRIANGLE.rotation=270
        _set_shape_transparency(RIGHT_TRIANGLE,10000)

        #add RIGHT_TRIANGLE
        left = util.Cm(0)
        top = util.Cm(13.05)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        RIGHT_TRIANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE_1.fill.solid()
        RIGHT_TRIANGLE_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #RIGHT_TRIANGLE_1.line
        RIGHT_TRIANGLE_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add RIGHT_TRIANGLE
        left = util.Cm(0)
        top = util.Cm(15.05)
        width = util.Cm(4)
        height = util.Cm(4)
        #shapes = slide.shapes
        RIGHT_TRIANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE_2.fill.solid()
        RIGHT_TRIANGLE_2.fill.fore_color.rgb = RGBColor(min(sc1+20, 255), min(sc2+20, 255), min(sc3+20, 255))
        #RIGHT_TRIANGLE_2.line
        RIGHT_TRIANGLE_2.line.color.rgb = RGBColor(min(sc1+20, 255), min(sc2+20, 255), min(sc3+20, 255))

        #add line
        start_x = util.Cm(2.95) 
        start_y = util.Cm(3.8)
        end_x = util.Cm(4.95)
        end_y = util.Cm(3.8)

        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
        line1.line.fill.background()
        line1.line.fill.solid()
        line1.line.fill.fore_color.rgb = RGBColor(0, 0, 0) #黑色

        #add title 
        left = util.Cm(2.95) 
        top = util.Cm(4.3)
        width = util.Cm(9.5)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add RECTANGL_1
        left = util.Cm(13.25)
        top = util.Cm(5.5)
        width = util.Cm(0.5)
        height = util.Cm(0.5)
        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #RIGHT_TRIANGLE_2.line
        RECTANGLE_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add RECTANGL2
        left = util.Cm(13.25)
        top = util.Cm(12.5)
        width = util.Cm(0.5)
        height = util.Cm(0.5)
        #shapes = slide.shapes
        RECTANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_2.fill.solid()
        RECTANGLE_2.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #RIGHT_TRIANGLE_2.line
        RECTANGLE_2.line.color.rgb = RGBColor(sc1, sc2, sc3)


        #add content1
        left = util.Cm(14.5) 
        top = util.Cm(3)
        width = util.Cm(9)
        height = util.Cm(6)
        #shapes = slide.shapes
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        # p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add content2
        left = util.Cm(14.5) 
        top = util.Cm(11)
        width = util.Cm(9)
        height = util.Cm(6)
        #shapes = slide.shapes
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        # p2.line_spacing = 1.5 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '7':
        #第九頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b ) 
        #add rectangle
        left = util.Cm(3.2)
        top = util.Cm(-3.19)
        width = util.Cm(18.9)
        height = util.Cm(25.6)
        #shapes = slide.shapes
        RIGHT_TRIANGLE = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE.fill.solid()
        RIGHT_TRIANGLE.fill.fore_color.rgb = RGBColor(224,224,224)
        #RIGHT_TRIANGLE_1.line
        #RIGHT_TRIANGLE.line.color.rgb = RGBColor(224,224,244)

        RIGHT_TRIANGLE.line.fill.background()

        RIGHT_TRIANGLE.rotation=270
        _set_shape_transparency(RIGHT_TRIANGLE,10000)

        #add RIGHT_TRIANGLE
        left = util.Cm(0)
        top = util.Cm(13.05)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        RIGHT_TRIANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE_1.fill.solid()
        RIGHT_TRIANGLE_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #RIGHT_TRIANGLE_1.line
        RIGHT_TRIANGLE_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add RIGHT_TRIANGLE
        left = util.Cm(0)
        top = util.Cm(15.05)
        width = util.Cm(4)
        height = util.Cm(4)
        #shapes = slide.shapes
        RIGHT_TRIANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE_2.fill.solid()
        RIGHT_TRIANGLE_2.fill.fore_color.rgb = RGBColor(min(sc1+20, 255), min(sc2+20, 255), min(sc3+20, 255))
        #RIGHT_TRIANGLE_2.line
        RIGHT_TRIANGLE_2.line.color.rgb = RGBColor(min(sc1+20, 255), min(sc2+20, 255), min(sc3+20, 255))

        #add line
        #add line
        start_x = util.Cm(2.95) 
        start_y = util.Cm(3.8)
        end_x = util.Cm(4.95)
        end_y = util.Cm(3.8)

        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
        line1.line.fill.background()
        line1.line.fill.solid()
        line1.line.fill.fore_color.rgb = RGBColor(0, 0, 0) #黑色

        #add title 
        left = util.Cm(2.95) 
        top = util.Cm(4.3)
        width = util.Cm(9)
        height = util.Cm(3)
        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add content
        left = util.Cm(2.95) 
        top = util.Cm(8)
        width = util.Cm(9)
        height = util.Cm(8)
        #shapes = slide.shapes
        RECTANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_2.fill.solid()
        RECTANGLE_2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_2.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0 
        p.text = f'{contents_done[0]}'

        RECTANGLE_2.line.fill.background()

        _set_shape_transparency(RECTANGLE_2,0)

        #add ROUNDRECTANGLE
        left = util.Cm(12.55)
        top = util.Cm(3.55)
        width = util.Cm(11)
        height = util.Cm(11.5)
        #shapes = slide.shapes
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(241,237,233)
        #rectangle.line
        rectangle.line.color.rgb = RGBColor(241,237,233)

        tf = rectangle.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note
        
    if pptchoose == '8':
        #第十頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背警
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b) 

            
        #add rectangle
        left = util.Cm(3) 
        top = util.Cm(3)
        width = util.Cm(15.5)
        height = util.Cm(12.5)
        #shapes = slide.shapes  #rectangle
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(241,237,233)
        rectangle_1.line.color.rgb = RGBColor(241,237,233)
            
        #add title
        left = util.Cm(4.5) 
        top = util.Cm(4.5)
        width = util.Cm(10)
        height = util.Cm(3)

        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

                    
        #add content
        left = util.Cm(4.5) 
        top = util.Cm(7.5)
        width = util.Cm(10)
        height = util.Cm(6)
        shapes = slide.shapes
        
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.2
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #加圖片
        #add ROUNDRECTANGLE
        left = util.Cm(15.5) 
        top = util.Cm(5.5)
        width = util.Cm(7.5)
        height = util.Cm(7.5)
        #shapes = slide.shapes
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(240,240,240)
        #rectangle.line
        rectangle.line.color.rgb = RGBColor(240,240,240)

        tf = rectangle.text_frame
        p = tf.paragraphs[0]
        p.text = '圖片'
        p.font.color.rgb=RGBColor(0,0,0)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = False #加粗 
        p.alignment = PP_ALIGN.CENTER
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '9':
        #第十一頁
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #add rectangle1
        left = util.Cm(0) 
        top = util.Cm(0)
        width = util.Cm(25.38)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle1.fill.solid()
        rectangle1.fill.fore_color.rgb = RGBColor(248,248,248)
        #rectangle.line
        rectangle1.line.fill.background()

        _set_shape_transparency(rectangle1,10000)

        #add rectangle2
        left = util.Cm(0)
        top = util.Cm(0)
        width = util.Cm(6)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        rectangle_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_2.fill.solid()
        rectangle_2.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_2.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add rectangle3
        left = util.Cm(3.5)
        top = util.Cm(8)
        width = util.Cm(9.5)
        height = util.Cm(8)
        #shapes = slide.shapes
        rectangle_3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_3.fill.solid()
        rectangle_3.fill.fore_color.rgb = RGBColor(stc1, stc2, stc3)
        #rectangle.line
        rectangle_3.line.color.rgb = RGBColor(stc1, stc2, stc3)

        #add rectangle4
        left = util.Cm(3)
        top = util.Cm(7.5)
        width = util.Cm(9.5)
        height = util.Cm(8)
        #shapes = slide.shapes
        rectangle_4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_4.fill.solid()
        rectangle_4.fill.fore_color.rgb = RGBColor(253,253,253)
        #rectangle.line
        rectangle_4.line.color.rgb = RGBColor(253,253,253)

        #add rectangle5
        left = util.Cm(14.5)
        top = util.Cm(8)
        width = util.Cm(9.5)
        height = util.Cm(8)
        #shapes = slide.shapes
        rectangle_5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_5.fill.solid()
        rectangle_5.fill.fore_color.rgb = RGBColor(stc1, stc2, stc3)
        #rectangle.line
        rectangle_5.line.color.rgb = RGBColor(stc1, stc2, stc3)

        #add rectangle6
        left = util.Cm(14)
        top = util.Cm(7.5)
        width = util.Cm(9.5)
        height = util.Cm(8)
        #shapes = slide.shapes
        rectangle_6 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_6.fill.solid()
        rectangle_6.fill.fore_color.rgb = RGBColor(253,253,253)
        #rectangle.line
        rectangle_6.line.color.rgb = RGBColor(253,253,253)

        #add title
        left = util.Cm(8.5) 
        top = util.Cm(3.5)
        width = util.Cm(10)
        height = util.Cm(3)

        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(20) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.CENTER
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add content1
        left = util.Cm(4) 
        top = util.Cm(8.5)
        width = util.Cm(7.5)
        height = util.Cm(6)
        shapes = slide.shapes
        
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(16) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add content2
        left = util.Cm(15) 
        top = util.Cm(8.5)
        width = util.Cm(7.5)
        height = util.Cm(6)
        shapes = slide.shapes
        
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(16) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '10':
        #第十二頁
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1

        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b ) 
        #add rectangle
        left = util.Cm(3.2)
        top = util.Cm(-3.19)
        width = util.Cm(18.9)
        height = util.Cm(25.6)
        #shapes = slide.shapes
        RIGHT_TRIANGLE = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE.fill.solid()
        RIGHT_TRIANGLE.fill.fore_color.rgb = RGBColor(224,224,224)
        #RIGHT_TRIANGLE_1.line
        #RIGHT_TRIANGLE.line.color.rgb = RGBColor(224,224,244)

        RIGHT_TRIANGLE.line.fill.background()

        RIGHT_TRIANGLE.rotation=270
        _set_shape_transparency(RIGHT_TRIANGLE,10000)

        #add rectangle1
        left = util.Cm(2.2)
        top = util.Cm(2)
        width = util.Cm(21)
        height = util.Cm(15)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add title
        left = util.Cm(3.2) 
        top = util.Cm(3.5)
        width = util.Cm(8)
        height = util.Cm(5)

        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add content
        left = util.Cm(12.2) 
        top = util.Cm(3.4)
        width = util.Cm(9)
        height = util.Cm(12)
        shapes = slide.shapes
        
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(16) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        #add oval
        i1 = 0 
        for i1 in range(3):

            left = util.Cm(20+i1*1) 
            top = util.Cm(2.5)
            width = util.Cm(0.5)
            height = util.Cm(0.5)
        #shapes = slide.shapes
            Oval_1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            Oval_1.fill.solid()
            Oval_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
            Oval_1.line.color.rgb = RGBColor(253,253,253)

        #add line
        start_x = util.Cm(22.2) 
        start_y = util.Cm(4)
        end_x = util.Cm(22.2)
        end_y = util.Cm(16)

        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,start_x, start_y, end_x, end_y)
        line1.line.fill.background()
        line1.line.fill.solid()
        line1.line.fill.fore_color.rgb = RGBColor(253,253,253) #黑色

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note
    
    if pptchoose == '11':
        #第十三頁
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = util.Cm(0)
        top = util.Cm(9.3)
        width = util.Cm(25.38)
        height = util.Cm(0.5)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add ROUNDED_RECTANGLE1
        left = util.Cm(16.5)
        top = util.Cm(2)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        ROUNDED_RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        ROUNDED_RECTANGLE_1.fill.solid()
        ROUNDED_RECTANGLE_1.fill.fore_color.rgb = RGBColor(240,240,240)
        #ROUNDED_RECTANGLE.line
        ROUNDED_RECTANGLE_1.line.color.rgb = RGBColor(240,240,240)

        #貼圖片
        tf = ROUNDED_RECTANGLE_1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.CENTER
        p1.text = '圖片'
        image_path= f'{img_list[0]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        #Create a shadow
        shadow = ROUNDED_RECTANGLE_1.shadow
        shadow.inherit = True
        shadow.visible = False
        shadow.distance = Pt(10)
        shadow.shadow_type = 'outer'
        shadow.angle = 45
        shadow.blur_radius = Pt(6)
        shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
        shadow.transparency = '50'
        shadow.distance = Pt(6)
        ROUNDED_RECTANGLE_1.shadow.style = 'outer'

        #add ROUNDED_RECTANGLE2
        left = util.Cm(3)
        top = util.Cm(11.1)
        width = util.Cm(6)
        height = util.Cm(6)
        #shapes = slide.shapes
        ROUNDED_RECTANGLE_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        ROUNDED_RECTANGLE_2.fill.solid()
        ROUNDED_RECTANGLE_2.fill.fore_color.rgb = RGBColor(240,240,240)
        #ROUNDED_RECTANGLE.line
        ROUNDED_RECTANGLE_2.line.color.rgb = RGBColor(240,240,240)

        tf = ROUNDED_RECTANGLE_2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(20) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.CENTER
        p2.text = '圖片'
        image_path= f'{img_list[1]}'
        slide.shapes.add_picture(image_path, left, top, width, height)

        #Create a shadow
        shadow = ROUNDED_RECTANGLE_2.shadow
        shadow.inherit = True
        shadow.visible = False
        shadow.distance = Pt(10)
        shadow.shadow_type = 'outer'
        shadow.angle = 45
        shadow.blur_radius = Pt(6)
        shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
        shadow.transparency = '50'
        shadow.distance = Pt(6)
        ROUNDED_RECTANGLE_2.shadow.style = 'outer'


        #add title
        left = util.Cm(3) 
        top = util.Cm(1)
        width = util.Cm(12)
        height = util.Cm(2)

        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add content1
        left = util.Cm(3) 
        top = util.Cm(4)
        width = util.Cm(12)
        height = util.Cm(4.5)
        shapes = slide.shapes
        
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(18) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)


        #add content2
        left = util.Cm(11) 
        top = util.Cm(11)
        width = util.Cm(12)
        height = util.Cm(5.5)
        shapes = slide.shapes
        
        RECTANGLE2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE2.fill.solid()
        RECTANGLE2.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE2.text_frame
        p2 = tf.paragraphs[0]
        p2.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p2.font.size = Pt(18) #大小
        p2.font.name = 'Arial' #字體
        p2.font.bold = False #加粗 
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.0 
        p2.text = f'{contents_done[1]}'

        RECTANGLE2.line.fill.background()

        _set_shape_transparency(RECTANGLE2,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note

    if pptchoose == '12':
        #第十四頁模板
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #投影片版配1
        #背景
        slide_bkg = slide.background
        slide_bkg_fill =  slide_bkg.fill
        slide_bkg_fill.solid()
        slide_bkg_fill.fore_color.rgb = RGBColor(c_1_r, c_1_g, c_1_b ) 
        #add rectangle
        left = util.Cm(0)
        top = util.Cm(0)
        width = util.Cm(25.08)
        height = util.Cm(19.05)
        #shapes = slide.shapes
        RIGHT_TRIANGLE = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
        RIGHT_TRIANGLE.fill.solid()
        RIGHT_TRIANGLE.fill.fore_color.rgb = RGBColor(224,224,224)
        #RIGHT_TRIANGLE_1.line
        #RIGHT_TRIANGLE.line.color.rgb = RGBColor(224,224,244)

        RIGHT_TRIANGLE.line.fill.background()

        _set_shape_transparency(RIGHT_TRIANGLE,10000)

        #add rectangle1
        left = util.Cm(13.5)
        top = util.Cm(4)
        width = util.Cm(10.5)
        height = util.Cm(12)
        #shapes = slide.shapes
        rectangle_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_1.fill.solid()
        rectangle_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
        #rectangle.line
        rectangle_1.line.color.rgb = RGBColor(sc1, sc2, sc3)

        #add rectangle2
        left = util.Cm(13)
        top = util.Cm(3.5)
        width = util.Cm(10.5)
        height = util.Cm(12)
        #shapes = slide.shapes
        rectangle_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_2.fill.solid()
        rectangle_2.fill.fore_color.rgb = RGBColor(253,253,253)
        #rectangle.line
        rectangle_2.line.color.rgb = RGBColor(253,253,235)

        #add oval
        i1 = 0 
        for i1 in range(3):

            left = util.Cm(21.5+i1*1) 
            top = util.Cm(2.5)
            width = util.Cm(0.5)
            height = util.Cm(0.5)
        #shapes = slide.shapes
            Oval_1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            Oval_1.fill.solid()
            Oval_1.fill.fore_color.rgb = RGBColor(sc1, sc2, sc3)
            Oval_1.line.color.rgb = RGBColor(253,253,253)
            
        #add title1
        left = util.Cm(2) 
        top = util.Cm(3.5)
        width = util.Cm(10)
        height = util.Cm(3)

        #shapes = slide.shapes
        RECTANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE_1.fill.solid()
        RECTANGLE_1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE_1.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p.font.size = Pt(28) #大小
        p.font.name = 'Arial' #字體
        p.font.bold = True #加粗 
        p.alignment = PP_ALIGN.LEFT
        p.text = f'{spl_copy[0]}'

        RECTANGLE_1.line.fill.background()

        _set_shape_transparency(RECTANGLE_1,0)

        #add content
        left = util.Cm(13.55) 
        top = util.Cm(4.3)
        width = util.Cm(9.3)
        height = util.Cm(10.5)
        shapes = slide.shapes
        
        RECTANGLE1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        RECTANGLE1.fill.solid()
        RECTANGLE1.fill.fore_color.rgb = RGBColor(253,253,253)

        tf = RECTANGLE1.text_frame
        p1 = tf.paragraphs[0]
        p1.font.color.rgb=RGBColor(stc1, stc2, stc3)
        p1.font.size = Pt(20) #大小
        p1.font.name = 'Arial' #字體
        p1.font.bold = False #加粗 
        p1.alignment = PP_ALIGN.LEFT
        p1.line_spacing = 1.0 
        p1.text = f'{contents_done[0]}'

        RECTANGLE1.line.fill.background()

        _set_shape_transparency(RECTANGLE1,0)

        text_note = slide.notes_slide.notes_text_frame
        text_note.text = note


        
    spl_copy.pop(0)
    print(len(spl_copy))
    try1.clear()
    img_list.clear()
    contents_done.clear()

def finish_ppt():
    global prs, con_word
    
    #第十五頁(結尾頁)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    #add RIGHT_TRIANGLE
    left = util.Cm(0)
    top = util.Cm(13.05)
    width = util.Cm(25.38)
    height = util.Cm(6)
    #shapes = slide.shapes
    RIGHT_TRIANGLE_1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
    RIGHT_TRIANGLE_1.fill.solid()
    RIGHT_TRIANGLE_1.fill.fore_color.rgb = RGBColor(mc1, mc2, mc3)
    #RIGHT_TRIANGLE_1.line
    RIGHT_TRIANGLE_1.line.color.rgb = RGBColor(mc1, mc2, mc3)

    #add text (end)
    left = width = util.Cm(8.5)
    top = util.Cm(8)
    height = util.Cm(3)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame


    p = tf.paragraphs[0]
    p.font.color.rgb=RGBColor(mtc1, mtc2, mtc3)
    p.font.size = Pt(48) #大小
    p.font.name = 'Arial' #字體
    p.font.bold = True #加粗 
    p.alignment = PP_ALIGN.CENTER
    p.text = 'The  End'
    
    ran = random.randint(10, 100000)
    prs.save(r'python_ppt_{:d}.pptx'.format(ran))
    print('done')


class SampleApp(tk.Tk):
    def __init__(self):
        tk.Tk.__init__(self)
        self._mainCanvas= None
        # The dictionary to hold the class type to switch to
        # Each new class passed here, will only have instance or object associated with it (i.e the result of the Key)
        self._allCanvases = dict()
        # Switch (and create) the single instance of StartUpPage
        self.switch_Canvas(homePage)

    def switch_Canvas(self, Canvas_class):

        # Unless the dictionary is empty, hide the current Frame (_mainCanvas is a frame)
        if self._mainCanvas:
            self._mainCanvas.pack_forget()

        # is the Class type passed one we have seen before?
        canvas = self._allCanvases.get(Canvas_class, False)

        # if Canvas_class is a new class type, canvas is False
        if not canvas:
            # Instantiate the new class
            canvas = Canvas_class(self)
            # Store it's type in the dictionary
            self._allCanvases[Canvas_class] = canvas

        # Pack the canvas or self._mainCanvas (these are all frames)
        canvas.pack(pady = 60)
        # and make it the 'default' or current one.
        self._mainCanvas = canvas

class homePage(tk.Frame):

    def __init__(self, master, *args, **kwargs):
        tk.Frame.__init__(self, master, *args, **kwargs)

        def save_topic():
            global key_word

            key_word = self.entry.get()
            print('topic: ', key_word)

        self.canvas = tk.Canvas(
            self,
            bg='#E26565', 
            height = 600,
            width = 900,
            bd = 0,
            highlightthickness = 0,
            relief = "ridge"
        )
        self.canvas.place(x = 0, y = 0)

        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            600.0,
            fill="#E8E9E9",
            outline=""
        )
        
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            50.0,
            fill="#8696A9",
            outline=""
        )
        
        self.canvas.create_text(
            350.0,
            10.0,
            anchor="nw",
            text="人工智慧簡報生成系統",
            fill="#FFFFFF",
            font=("InriaSans Regular", 20 * -1,'bold')
        )

        self.canvas.create_text(
            15.0,
            59.0,
            anchor="nw",
            text="歡迎您使用人工智慧簡報生成系統，讓您輕鬆生成一份完美的簡報，讓我們開始這趟神奇的旅程吧~",
            fill="#000000",
            font=("InriaSans Regular", 20 * -1)
        )
        
        self.button = ctk.CTkButton(
            self,
            text = "開始",
            text_color= 'white', #文字顏色
            command=lambda: [master.switch_Canvas(PageOne),save_topic()],
            fg_color='#8696A9', #按鈕顏色
            corner_radius= 30, #邊角圓弧度
            hover_color='#7A8797',
            bg_color="#E8E9E9",
            font=('Helvetica', 25, 'bold'),
            width=150.0,
            height=63.0
        )
        
        self.button.place(
            x=375.0,
            y=355.0,
        )

        self.entry = tk.Entry(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            font=('Aries', 16),
            highlightthickness=0
        )

        self.entry.place(
            x=181.0,
            y=225.0,
            width=539.0,
            height=45.0
        )

        self.canvas.create_text(
            181.0,
            187.0,
            anchor="nw",
            text="今天想做的簡報主題:",
            fill="#000000",
            font=("Inter Bold", 24 * -1)
        )
        
        self.buttonr = ctk.CTkButton(
            self,
            command=lambda: finish_ppt(),
            fg_color='red', #按鈕顏色
            text="",
            # corner_radius= 30, #邊角圓弧度
            # hover_color='#7A8797',
            bg_color="red",
            # font=('Helvetica', 25, 'bold'),
            width=20.0,
            height=10.0
        )
        
        self.buttonr.place(
            x=880.0,
            y=590.0,
        )
        
        # self.button.pack(side="top", fill="x", pady=5)
        self.canvas.pack()
        # self.canvas.pack(fill=tk.BOTH, side=tk.LEFT, expand=True)
    
class PageOne(tk.Frame):  #chatGPT

    def __init__(self, master, *args, **kwargs):
        tk.Frame.__init__(self,master, *args, **kwargs)
        
        #--chatGPT - generate titles
        def generate_titles():
            global spl,chatans,key_word,tousr, amounts, spl_copy, regenerate, do

            chatans.clear()
            self.entry1.delete('1.0', tk.END)
            spp = ''
            amounts = self.entry3.get()
            #     print(key_words)
            prompts = f'我想製作一個以"{key_word}"為主題的簡報，總共需要{amounts}個與"{key_word}"專業相關的標題，請幫我用繁體中文生成，標題與標題之間最好有連貫性、有邏輯、有多元性，但不可以重複。'
            if regenerate == True :
                prompts = '你剛剛生成的我覺得不太適合，請重新生成。' + prompts
            regenerate = True
            doQA(prompts, 512, 0.2+do)
            spp = ''.join(chatans[0])
            spl = spp.split('\n')
            tousr = ""

            

            for i in range(10):
                if len(spl[0])< 5:
                    print(spl[0])
                    spl.pop(0)
                
            for i in range(len(spl)):
                print(f"the word:{spl[i]}")
                spl[i] = spl[i][2:]
                spl[i] = spl[i].lstrip('.')

            spl = [i.strip() for i in spl if i.strip() != '']
            # print('spl: ', spl)
            # spl_copy = spl.copy()
            # print('spl_copy: ', spl_copy)
            
            for i in range(0, len(spl)):
                if i == 0:
                    self.entry1.insert('1.0', f'{i + 1}. {spl[i]}\n')
                    
                elif i < len(spl) - 1 and i > 0:
                    self.entry1.insert('end', f'{i + 1}. {spl[i]}\n')
                
                else:
                    self.entry1.insert('end', f'{i + 1}. {spl[i]}')
            print("spl(before): ", spl)
            
        def give_hint():
            global hint_yn
    
            if hint_yn == 0:
                self.canvas.delete('origin')
                self.canvas.create_text(
                    20.0,
                    10.0,
                    anchor="nw",
                    text="您可以手動調整標題內容，確認好後請點擊「下一步」",
                    fill="#FFFFFF",
                    font=("InriaSans Regular", 20 * -1),
                    tags= 'hint'
                )
                hint_yn = 1

        def hint_to_0():
            global hint_yn
            hint_yn = 0
            
        def save_spl():
            global spl, spl_copy
            spl.clear()
            spl_copy.clear()
            spl_cont = self.entry1.get(1.0, "end")
            spl_split = spl_cont.split('\n')
            spl_split.pop(len(spl_split) - 1)
            for i in range(len(spl_split)):
                spl_split[i] = spl_split[i].lstrip(f'{i + 1}')
                spl_split[i] = spl_split[i].lstrip('.')
                spl_split[i] = spl_split[i].lstrip(' ')
                
            spl = spl_split.copy()
            spl_copy = spl_split.copy()
            print("spl: ", spl, '\nspl_copy: ', spl_copy)
            
        def tocolortp():
            global mc1, mc2, mc3, mcname
            colors = askcolor(title="Tkinter Color Chooser")
            print("r: ", colors[0][0]) # r
            mc1 = colors[0][0]
            print("g: ", colors[0][1]) # g
            mc2 = colors[0][1]
            print("b: ", colors[0][2]) # b
            mc3 = colors[0][2]
            mcname = colors[1]
            showTcolor()
            
        def tocolortpt():
            global mtc1, mtc2, mtc3, mtcname
            colors = askcolor(title="Tkinter Color Chooser")
            print("r: ", colors[0][0]) # r
            mtc1 = colors[0][0]
            print("g: ", colors[0][1]) # g
            mtc2 = colors[0][1]
            print("b: ", colors[0][2]) # b
            mtc3 = colors[0][2]
            mtcname = colors[1]
            showTtcolor()


        def showTcolor():
            global mc1, mc2, mc3, mcname
            self.Tbutton.configure(bg = f'{mcname}')
        
        def showTtcolor():
            global mtc1, mtc2, mtc3, mtcname
            self.Ttbutton.configure(bg = f'{mtcname}')

        self.canvas = tk.Canvas(
            self,
            bg = "#E26565",
            height = 600,
            width = 900,
            bd = 0,
            highlightthickness = 0,
            relief = "ridge"
        )
        
        self.buttonr = ctk.CTkButton(
            self,
            command=lambda: finish_ppt(),
            fg_color='red', #按鈕顏色
            text="",
            # corner_radius= 30, #邊角圓弧度
            # hover_color='#7A8797',
            bg_color="red",
            # font=('Helvetica', 25, 'bold'),
            width=20.0,
            height=10.0
        )
        
        self.buttonr.place(
            x=880.0,
            y=590.0,
        )

        self.canvas.place(x = 0, y = 0)
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            600.0,
            fill="#E8E9E9",
            outline=""
        )

        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            50.0,
            fill="#8696A9",
            outline=""
        )

        self.canvas.create_text(
            20.0,
            10.0,
            anchor="nw",
            text="輸入欲生成的子標題數後，請點擊「生成標題」來產生簡報的子標題",
            fill="#FFFFFF",
            font=("InriaSans Regular", 20 * -1),
            tags= 'origin'
        )

#下一頁按鈕
        self.button2 = ctk.CTkButton(
            self,
            text = "下一步",
            text_color= 'white', #文字顏色
            fg_color="#8696AF", #按鈕顏色
            hover_color='#7A8797',
            bg_color="#E8E9E9",
            border_color="#8696AF",
            command=lambda: [save_spl(), master.switch_Canvas(PageTwo),generate_ppt(),hint_to_0(), change_regenerate(), temperature_0()],
            font=('Aries', 20),
            corner_radius= 0, #邊角圓弧度
            width=98.0,
            height=50.0
        )

        self.button2.place(
            x=802.0,
            y=0.0
        )

        self.entry1 = tk.Text(
            self,
            padx=15,
            pady=5,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            highlightthickness=0,
            font=('Aries', 14)
        )

        self.entry1.place(
            x=32.0,
            y=140.0,
            width=560.0,
            height=430.0
        )

        self.canvas.create_text(
            32.0,
            73.0,
            anchor="nw",
            text="欲生成的子標題數: ",
            fill="#000000",
            font=("Inter Bold", 20 * -1)
        )
        
        self.canvas.create_text(
            32.0,
            110.0,
            anchor="nw",
            text="生成的標題：",
            fill="#000000",
            font=("Inter Bold", 20 * -1)
        )

        self.entry3 = tk.Entry(
            self,
            bg="#FFFFFF",
            fg="#000716",
            relief="ridge",
            justify='center',
            font=('Aries', 16, 'bold'),
            highlightthickness=0
        )

        self.entry3.place(
            x=205.0,
            y=62.0,
            width=50.0,
            height=45.0
        )

        self.button3 = ctk.CTkButton(
            self,
            text = "生成標題",
            text_color= 'white', #文字顏色
            bg_color="#E8E9E9", 
            command=lambda: [generate_titles(),give_hint(),temperature_up()],
            fg_color='#8696A9', #按鈕顏色
            corner_radius= 50, #邊角圓弧度
            hover_color='#7A8797', 
            font=('Helvetica', 20),
            width=150.0,
            height=35.0
        )

        self.button3.place(
            x=270.0,
            y=68.0
        )
        
        self.textLC = self.canvas.create_text(
            660.0,
            120.0,
            anchor="nw",
            text=f"封面/目錄模板顏色: ",
            fill="#000000",
            font=("InriaSans Bold", 12)
        )   
        
        self.buttonLC = ctk.CTkButton(
            self,
            text = "選擇",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: tocolortp(),
            font=('InriaSans Bold', 16),
            corner_radius= 30, #邊角圓弧度
            width=90.0,
            height=30.0
        )

        self.buttonLC.place(
            x=800.0,
            y=112.0
        )
        
        self.Tbutton = tk.Label(
            self,
            bg = '#705E51',
            width =  30,
            height = 2,
            justify = 'center',
        )
        
        self.Tbutton.place(
            x = 660,
            y = 150
        )
        
        self.textTC = self.canvas.create_text(
            660.0,
            200.0,
            anchor="nw",
            text=f"封面/目錄字體顏色: ",
            fill="#000000",
            font=("InriaSans Bold", 12)
        ) 
        
        self.buttonTC = ctk.CTkButton(
            self,
            text = "選擇",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: tocolortpt(),
            font=('InriaSans Bold', 16),
            corner_radius= 30, #邊角圓弧度
            width=90.0,
            height=30.0
        )

        self.buttonTC.place(
            x=800.0,
            y=195.0
        )
        
        self.Ttbutton = tk.Label(
            self,
            bg = '#504C50',
            width =  30,
            height = 2,
            justify = 'center',
        )
        
        self.Ttbutton.place(
            x = 660,
            y = 230
        )

        self.canvas.pack()

class PageTwo(tk.Frame):  #chatGPT

    def __init__(self, master, *args, **kwargs):
        tk.Frame.__init__(self,master, *args, **kwargs)
        global page_of_ppt, layout, picpath, pagechangei, pindex, spl, spl_copy
        # print("spl(p2): ", spl)
        # print("page"+r':{:d}'.format(page_of_ppt))
        
        def clearContents():
            global page_of_ppt, spl
            print('page_of_ppt: ', page_of_ppt)
            self.entry3.delete('-1', tk.END)
            if page_of_ppt+1 < len(spl):
                self.entry3.insert(END, spl[page_of_ppt+1])

        def changePPTpage():
            global page_of_pic, pic_pic, pic_path, layout, pindex, picpath, pagechange, pagechangei
            picpath = 0
            
            if (layout + 3) < len(page_of_pic):
                pindex = layout
                for i in range(4):
                    pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
                    pic_path[picpath].config(image=pic_pic[layout])
                    picpath += 1
                    layout += 1
                
                ppt_page = [
                    ppt_page_index[pindex],
                    ppt_page_index[pindex + 1],
                    ppt_page_index[pindex + 2],
                    ppt_page_index[pindex + 3]
                ]
                self.combo.config(values = ppt_page)
                
                
            else:
                layout = 0
                pindex = layout
                for i in range(4):
                    pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
                    pic_path[picpath].config(image=pic_pic[layout])
                    picpath += 1
                    layout += 1
                    
                ppt_page = [
                    ppt_page_index[pindex],
                    ppt_page_index[pindex + 1],
                    ppt_page_index[pindex + 2],
                    ppt_page_index[pindex + 3]
                ]
                self.combo.config(values = ppt_page)
                

        def changePPTpageback():
            global page_of_pic, pic_pic, pic_path, layout, picpath, pagechange, pagechangei, pagechange1, pagechangei1
            
            if layout == 4:
                picpath = 0
                layout = 8
                for i in range(4):
                    pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
                    pic_path[picpath].config(image=pic_pic[layout])
                    picpath += 1
                    layout += 1
                    
                ppt_page = [
                    ppt_page_index[8],
                    ppt_page_index[9],
                    ppt_page_index[10],
                    ppt_page_index[11]
                ]
                self.combo.config(values = ppt_page)
                
            elif layout == 8:
                layout = 0
                picpath = 0
                for i in range(4):
                    pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
                    pic_path[picpath].config(image=pic_pic[layout])
                    picpath += 1
                    layout += 1
                    
                ppt_page = [
                    ppt_page_index[0],
                    ppt_page_index[1],
                    ppt_page_index[2],
                    ppt_page_index[3]
                ]
                self.combo.config(values = ppt_page)
                
            elif layout == 12:
                layout = 4
                picpath = 0
                for i in range(4):
                    pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
                    pic_path[picpath].config(image=pic_pic[layout])
                    picpath += 1
                    layout += 1
                    
                ppt_page = [
                    ppt_page_index[4],
                    ppt_page_index[5],
                    ppt_page_index[6],
                    ppt_page_index[7]
                ]
                self.combo.config(values = ppt_page)
                
                   
        def check_choose_page():
            global pptchoose, cont, picc
            pptchoose = self.combo.get()
            print('pptchoose: ', pptchoose, type(pptchoose))
            if pptchoose == '1' or pptchoose == '2' or pptchoose == '7' or pptchoose == '8':
                cont = 1
                picc = 1
            elif pptchoose == '4' or pptchoose == '6' or pptchoose == '9':
                cont = 2
                picc = 0
            elif pptchoose == '3' or pptchoose == '11':
                cont = 2
                picc = 2
            elif pptchoose == '5':
                cont = 3
                picc = 0
            elif pptchoose == '10' or pptchoose == '12':
                cont = 1
                picc = 0
                
            print('cont: ', cont, '\npicc: ', picc)
            
        def tocolortp():
            global sc1, sc2, sc3, scname
            colors = askcolor(title="Tkinter Color Chooser")
            print("r: ", colors[0][0]) # r
            sc1 = colors[0][0]
            print("g: ", colors[0][1]) # g
            sc2 = colors[0][1]
            print("b: ", colors[0][2]) # b
            sc3 = colors[0][2]
            scname = colors[1]
            showTcolor()
            
        def tocolortpt():
            global stc1, stc2, stc3, stcname
            colors = askcolor(title="Tkinter Color Chooser")
            print("r: ", colors[0][0]) # r
            stc1 = colors[0][0]
            print("g: ", colors[0][1]) # g
            stc2 = colors[0][1]
            print("b: ", colors[0][2]) # b
            stc3 = colors[0][2]
            stcname = colors[1]
            showTtcolor()
        
        def showTcolor():
            global mc1, mc2, mc3, scname
            self.Tbutton.configure(bg = f'{scname}')
        
        def showTtcolor():
            global mtc1, mtc2, mtc3, stcname
            self.Ttbutton.configure(bg = f'{stcname}')

#視窗大小
        self.canvas = tk.Canvas(
            self,
            bg = "#E26565",
            height = 600,
            width = 900,
            bd = 0,
            highlightthickness = 0,
            relief = "ridge"
        )

        self.canvas.place(x = 0, y = 0)
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            600.0,
            fill="#E8E9E9",
            outline=""
        )
#介面上方padding
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            50.0,
            fill="#8696A9",
            outline=""
        )

        self.canvas.create_text(
            20.0,
            10.0,
            anchor="nw",
            text="選擇你這一頁想要的簡報模板，選好後請點擊「下一步」繼續",
            fill="#FFFFFF",
            font=("InriaSans Regular", 20 * -1)
        )

#下一頁按鈕
        self.button2 = ctk.CTkButton(
            self,
            text = "下一步",
            text_color= 'white', #文字顏色
            fg_color="#8696AF", #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            border_color="#8696AF",
            command=lambda: [clearContents(), master.switch_Canvas(temp), check_choose_page()],
            font=('Aries', 20),
            corner_radius= 0, #邊角圓弧度
            width=98.0,
            height=50.0
        )

        self.button2.place(
            x=802.0,
            y=0.0
        )
        
        self.buttonr = ctk.CTkButton(
            self,
            command=lambda: finish_ppt(),
            fg_color='red', #按鈕顏色
            text="",
            # corner_radius= 30, #邊角圓弧度
            # hover_color='#7A8797',
            bg_color="red",
            # font=('Helvetica', 25, 'bold'),
            width=20.0,
            height=10.0
        )
        
        self.buttonr.place(
            x=880.0,
            y=590.0,
        )
        
        self.canvas.create_text(
            32.0,
            73.0,
            anchor="nw",
            text="此頁標題:",
            fill="#000000",
            font=("Inter Bold", 20 * -1)
        )
        
        self.entry3 = tk.Entry(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            highlightthickness=0
        )

        self.entry3.insert(END, spl[page_of_ppt])
        
        self.entry3.place(
            x=150.0,
            y=62.0,
            width=162.0,
            height=45.0
        )

        ppt_page =[
            ppt_page_index[pindex],
            ppt_page_index[pindex + 1],
            ppt_page_index[pindex + 2],
            ppt_page_index[pindex + 3]
        ]
        
        self.combo = ttk.Combobox(
            self, 
            values = ppt_page, 
            state = 'readonly',
            width = 10
        )
        
        self.combo.place(
            x=350.0,
            y=70.0
        )
        
        index = 0
        for i in range(4):
            pic_pic[layout] = PhotoImage(file=f'{page_of_pic[layout]}')
            if i < 2:
                pic_path[picpath] = tk.Label(
                    self,
                    image=pic_pic[picpath],
                    width=302,
                    height=225
                )
                    
                pic_path[picpath].place(
                    x=20.0 + 330 * i,
                    y=120.0
                )
            else:
                pic_path[picpath] = tk.Label(
                    self,
                    image=pic_pic[picpath],
                    width=302,
                    height=225
                )
                    
                pic_path[picpath].place(
                    x=20.0 + 330 * index,
                    y=350.0
                )
                
                index += 1
            
            layout += 1
            picpath += 1
        
        self.textLC = self.canvas.create_text(
            660.0,
            120.0,
            anchor="nw",
            text=f"簡報模板顏色: ",
            fill="#000000",
            font=("InriaSans Bold", 12)
        )   
        
        self.buttonLC = ctk.CTkButton(
            self,
            text = "選擇",
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: tocolortp(),
            font=('InriaSans Bold', 16),
            corner_radius= 30, #邊角圓弧度
            width=90.0,
            height=30.0
        )

        self.buttonLC.place(
            x=800.0,
            y=112.0
        )
        
        self.Tbutton = tk.Label(
            self,
            bg = '#705E51',
            width =  30,
            height = 2,
            justify = 'center',
        )
        
        self.Tbutton.place(
            x = 660,
            y = 150
        )
        
        self.textTC = self.canvas.create_text(
            660.0,
            200.0,
            anchor="nw",
            text=f"簡報字體顏色: ",
            fill="#000000",
            font=("InriaSans Bold", 12)
        ) 
        
        self.buttonTC = ctk.CTkButton(
            self,
            text = "選擇",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: tocolortpt(),
            font=('InriaSans Bold', 16),
            corner_radius= 30, #邊角圓弧度
            width=90.0,
            height=30.0
        )

        self.buttonTC.place(
            x=800.0,
            y=195.0
        )
        
        self.Ttbutton = tk.Label(
            self,
            bg = '#504C50',
            width =  30,
            height = 2,
            justify = 'center',
        )
        
        self.Ttbutton.place(
            x = 660,
            y = 230
        )
        
        self.button6 = ctk.CTkButton(
            self,
            text = '模板上一頁',
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: changePPTpageback(),
            font=('Aries', 16),
            corner_radius= 30, #邊角圓弧度
            width=100.0,
            height=35.0
        )

        self.button6.place(
            x=480.0,
            y=70.0
        )
        
        self.button7 = ctk.CTkButton(
            self,
            text = '模板下一頁',
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: changePPTpage(),
            font=('Aries', 16),
            corner_radius= 30, #邊角圓弧度
            width=100.0,
            height=35.0
        )

        self.button7.place(
            x=600.0,
            y=70.0
        )
            
                
        self.canvas.pack()
        
class temp(tk.Frame):
    def __init__(self, master, *args, **kwargs):
        tk.Frame.__init__(self,master, *args, **kwargs)
        global page_of_ppt, spl, cont
        
        def finishornot():
            global spl, page_of_ppt
            if(page_of_ppt == len(spl)):
                return True
            else:
                return False
        
        def piccornot():
            global picc, contents_done, spl_copy
            if(picc != 0):
                master.switch_Canvas(PageThree)
            else:
                generate_one_page_ppt()
                if(finishornot() == True):
                    finish_ppt()
                    master.switch_Canvas(homePage)
                else:
                    master.switch_Canvas(PageTwo)
        
        def clearContents():
            global page_of_ppt, spl
            self.entry3.delete('-1', tk.END)
            self.entry2.delete('1.0', tk.END)
            self.entry1.delete('1.0', tk.END)
            self.entry4.delete('1.0', tk.END)
            if page_of_ppt < len(spl):
                self.entry3.insert(END, spl[page_of_ppt])
            
        def check_contents():
            global spp2,spl2,try1,spl, tindex, contents_done, spp3, cont, spl_copy, consave, note
                
            contents_done.clear()
            spl2.clear()
            contenttt = self.entry1.get(1.0, "end")
            note = self.entry4.get(1.0, "end")
            print("content: ", contenttt)

            for i in range(1,len(spl)):
                spl[i] = spl[i].replace(' ',"")
            # print('spl: ', spl)
            spl_copy[0] = self.entry3.get()
            # print(spl_copy[0])
            
            if cont > 1:
                spp3.clear()
                spp3 = contenttt.split('\n')
                print(len(spp3))
                for i in range(len(spp3)):
                    for j in range(3):
                        spp3[i] = clean(spp3[i], False)
                    if len(spp3[i]) > 2:
                        contents_done.append(spp3[i])
                        print(f"第{i}"+spp3[i])
                print('contents_done: ', contents_done)
                
            elif cont == 1:
                contenttt = contenttt.replace(' ', '')
                contenttt = contenttt.replace('\n', '')
                contents_done.append(contenttt)
                print('contents_done: ', contents_done)    
            
        def generate_contents():
            global chatans,tousr,spp2,page_of_ppt, cont, spp3, note, regenerate, do
            note = ''
            spp2 = ''
            spp3.clear()
            chatans.clear()
            self.entry1.delete('1.0', tk.END)
            self.entry4.delete('1.0', tk.END)
            ttopic = self.entry3.get()
            if cont == 1:
                if self.combo.get() == "依我提供的資料擴充生成內文":
                    word = self.entry2.get(1.0, "end")
                    prompts = f'''請你針對"{ttopic}"這個標題與"{word}"這段話，擴充生成適合的內文，總共一句。'''
                    if regenerate == True :
                        prompts = '你剛剛生成的我覺得不太適合，請重新生成。' + prompts
                    regenerate = True
                    doQA(prompts, 2048, 0.2+do)
                    spp2 = ''.join(chatans[0])
                    print(spp2)
                    spp2 = clean(spp2, True)
                    spp3 = spp2.split('\n')
                    spp2 = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            spp2 = spp3[i]
                    self.entry1.insert('1.0', spp2)
                    prompts = f'''請你幫"{spp2}"段文字擴充生成專業的詳細的介紹，總共一段。'''
                    spp3.clear()
                    chatans.clear()
                    doQA(prompts, 2048, 0.2+do)
                    note = ''.join(chatans[0])
                    print("講稿修改前為:"+note)
                    note = clean(note, True)
                    spp3 = note.split('\n')
                    note = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            note = note + spp3[i]
                    self.entry4.insert('1.0', note)
                else:
                    if self.combo.get() == "依我提供的資料縮減生成內文":
                        word = self.entry2.get(1.0, "end")
                        print("我接收到的文字:" + word)
                        prompts = f'''請用繁體中文摘錄"{word}"這段文字的重點，總共一段'''
                        print("pr" + prompts)
                    elif self.combo.get() == "生成全新的內文":
                        prompts = f'''請你針對"{ttopic}"這個標題，提供值得學術界探討的專業的相關重要議題。麻煩請用繁體中文生成100字左右的敘述，總共一段，中間請勿分段'''
                    if regenerate == True :
                        prompts = '你剛剛生成的我覺得不太適合，請重新生成。' + prompts
                    regenerate = True
                    doQA(prompts, 2048, 0.2+do)
                    note = ''.join(chatans[0])
                    print(note)
                    note = clean(note, True)
                    spp3 = note.split('\n')
                    note = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('講稿修改後為:'+spp3[i])
                            note = note + spp3[i]
                    self.entry4.insert('1.0', note)
                    prompts = f'''請從"{note}"這段文字摘取重點，總共一句，請不要提及"學術界"'''
                    spp3.clear()
                    chatans.clear()
                    doQA(prompts, 2048, 0.2+do)
                    spp2 = ''.join(chatans[0])
                    print(spp2)
                    spp2 = clean(spp2, True)
                    spp3 = spp2.split('\n')
                    spp2 = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            spp2 = spp2 + spp3[i]
                    self.entry1.insert('1.0', spp2)
                    
                
            elif cont > 1:
                if self.combo.get() == "依我提供的資料擴充生成內文":
                    word = self.entry2.get(1.0, "end")
                    prompts = f'''我要製作一頁有關"{ttopic}"的簡報，請幫我尋找有關"{word}"的相關資料，接著參考有關"{word}"的相關資料後用繁體中文幫我生成適合放進這個簡報的專業內文，總共生成{cont}句，每一句不可以重複。'''
                    if regenerate == True :
                        prompts = '你剛剛生成的我覺得不太適合，請重新生成。' + prompts
                    regenerate = True
                    doQA(prompts, 2048, 0.2+do)
                    spp2 = ''.join(chatans[0])
                    print(spp2)
                    spp2 = clean(spp2, True)
                    spp3 = spp2.split('\n')
                    spp2 = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            spp2 = spp2 + spp3[i] + '\n' + '\n'
                    self.entry1.insert('1.0', spp2)
                    prompts = f'''我要製作一頁有關"{ttopic}"的簡報，請提供一份包含以下提供的文句："{spp2}"且進行深入探討的專業又詳細的講稿，總共{cont}段。'''
                    spp3.clear()
                    chatans.clear()
                    doQA(prompts, 2048, 0.2+do)
                    note = ''.join(chatans[0])
                    print(note)
                    note = clean(note, True)
                    spp3 = note.split('\n')
                    note = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            note = note + spp3[i] + '\n' + '\n'
                    self.entry4.insert('1.0', note)

                else:
                    if self.combo.get() == "依我提供的資料縮減生成內文":
                        word = self.entry2.get(1.0, "end")
                        prompts = f'''我要製作一頁有關"{ttopic}"的簡報，請用繁體中文摘錄"{word}"這段文字的重點，總共{cont}段'''
                    elif self.combo.get() == "生成全新的內文":
                        prompts = f'''請你針對"{ttopic}"這個標題，提供值得學術界探討的專業的相關重要議題。麻煩請用繁體中文生成100字左右，總共{cont}段。然後生成對應的簡報內文講稿'''
                    if regenerate == True :
                        prompts = '你剛剛生成的我覺得不太適合，請重新生成。' + prompts
                    regenerate = True
                    doQA(prompts, 2048, 0.2+do)
                    note = ''.join(chatans[0])
                    print(note)
                    note = clean(note, True)
                    spp3 = note.split('\n')
                    note = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            note = note + spp3[i] + '\n' + '\n'
                    self.entry4.insert('1.0', note)
                    prompts = f'''我要製作一頁有關"{ttopic}"的簡報，請從"{note}"這段文字摘取重點，總共{cont}句'''
                    spp3.clear()
                    chatans.clear()
                    doQA(prompts, 2048, 0.2+do)
                    spp2 = ''.join(chatans[0])
                    print(spp2)
                    spp2 = clean(spp2, True)
                    spp3 = spp2.split('\n')
                    spp2 = ''
                    for i in range(len(spp3)):
                        if len(spp3[i])>6:
                            print('it'+spp3[i])
                            spp2 = spp2 + spp3[i] + '\n' + '\n'
                    self.entry1.insert('1.0', spp2)

        def clean(c, yn):
            if yn == True:
                for i in range(3): 
                    c = c.lstrip('。')
                    c = c.lstrip('，')
                    c = c.lstrip(':')
                    c = c.lstrip('：')
                    c = c.lstrip('')
                    c = c.lstrip('\n')
                    c = c.rstrip('')
                    c = c.rstrip('\n')
            else :
                for i in range(3):
                    c = c.lstrip('1')
                    c = c.lstrip('2')
                    c = c.lstrip('3')
                    c = c.lstrip('.')
                    c = c.lstrip('第一段')
                    c = c.lstrip('第二段')
                    c = c.lstrip('第三段')
                    c = c.lstrip('段落一')
                    c = c.lstrip('段落二')
                    c = c.lstrip('段落三')
                    c = c.lstrip('：')
                    c = c.lstrip(':')
                    c = c.lstrip(' ')
                    c = c.lstrip('\n')
                    c = c.lstrip('，')
                    c = c.lstrip('。')
                    c = c.lstrip('')
            return c

        def give_hint():
            global hint_yn
    
            if hint_yn == 0:
                self.canvas.delete('origin')
                
                self.canvas.create_text(
                    20.0,
                    10.0,
                    anchor="nw",
                    text="您可以手動調整內文內容，確認好後請點擊「確認內文」",
                    fill="#FFFFFF",
                    font=("InriaSans Regular", 20 * -1),
                    tags= 'hint1'
                )
                hint_yn = 1

            elif hint_yn == 1:
                self.canvas.delete('hint1')
                self.canvas.create_text(
                    20.0,
                    10.0,
                    anchor="nw",
                    text="接下來請按「下一步」，進入生成圖片或製作下一頁投影片",
                    fill="#FFFFFF",
                    font=("InriaSans Regular", 20 * -1),
                    tags= 'hint2'
                )
                hint_yn = 1


        def hint_to_0():
            global hint_yn
            hint_yn = 0
            self.canvas.delete('hint2')
            self.canvas.create_text(
                20.0,
                10.0,
                anchor="nw",
                text="您可以在選單中選擇生成內文的方式，選好後請點擊「生成內文」",
                fill="#FFFFFF",
                font=("InriaSans Regular", 20 * -1),
                tags= 'origin'
            )


        self.canvas = tk.Canvas(
            self,
            bg = "#E26565",
            height = 600,
            width = 900,
            bd = 0,
            highlightthickness = 0,
            relief = "ridge"
        )

        self.canvas.place(x = 0, y = 0)
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            600.0,
            fill="#E8E9E9",
            outline=""
        )

        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            50.0,
            fill="#8696A9",
            outline=""
        )

        self.canvas.create_text(
            20.0,
            10.0,
            anchor="nw",
            text="您可以在選單中選擇生成內文的方式，選好後請點擊「生成內文」",
            fill="#FFFFFF",
            font=("InriaSans Regular", 20 * -1),
            tags= 'origin'
        )
        
        self.button2 = ctk.CTkButton(
            self,
            text = "下一步",
            text_color= 'white', #文字顏色
            fg_color="#8696AF", #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            border_color="#8696AF",
            command=lambda: [next(), clearContents(), finishornot(), piccornot(),hint_to_0(), change_regenerate(), temperature_0()],
            font=('Aries', 20),
            corner_radius= 0, #邊角圓弧度
            width=98.0,
            height=50.0
        )

        self.button2.place(
            x=802.0,
            y=0.0
        )
        
        self.buttonr = ctk.CTkButton(
            self,
            command=lambda: finish_ppt(),
            fg_color='red', #按鈕顏色
            text="",
            # corner_radius= 30, #邊角圓弧度
            # hover_color='#7A8797',
            bg_color="red",
            # font=('Helvetica', 25, 'bold'),
            width=20.0,
            height=10.0
        )
        
        self.buttonr.place(
            x=880.0,
            y=590.0,
        )
        
        self.canvas.create_text(
            32.0,
            73.0,
            anchor="nw",
            text="此頁標題:",
            fill="#000000",
            font=("Inter Bold", 20 * -1)
        )
        
        self.entry3 = tk.Entry(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            highlightthickness=0
        )

        self.entry3.insert(END, spl[page_of_ppt])
        
        self.entry3.place(
            x=150.0,
            y=62.0,
            width=162.0,
            height=45.0
        )
        
        generate_content_kinds =[
            "生成全新的內文",
            "依我提供的資料縮減生成內文",
            "依我提供的資料擴充生成內文",
        ]
        
        self.combo = ttk.Combobox(
            self, 
            values = generate_content_kinds, 
            state = 'readonly',
            height= 30,
            width= 25
        )
        
        self.combo.place(
            x=325.0,
            y=70.0
        )

        self.combo.current(0)
#講稿output
        self.canvas.create_text(
            32.0,
            355.0,
            anchor="nw",
            text="生成的講稿:",
            fill="#000000",
            font=("Inter Bold", 16 * -1)
        )

        self.entry4 = tk.Text(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            font=("標楷體", 16 * -1),
            highlightthickness=0
        )

        self.entry4.place(
            x=25.0,
            y=377.0,
            width=850.0,
            height=200.0
        )
#input
        self.canvas.create_text(
            32.0,
            130.0,
            anchor="nw",
            text="您提供的文字:",
            fill="#000000",
            font=("Inter Bold", 16 * -1)
        )

        self.entry2 = tk.Text(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            font=("標楷體", 16 * -1),
            highlightthickness=0
        )

        self.entry2.place(
            x=25.0,
            y=150.0,
            width=315.0,
            height=200.0
        )

#內文output
        self.canvas.create_text(
            400.0,
            130.0,
            anchor="nw",
            text="生成的簡報內文:",
            fill="#000000",
            font=("Inter Bold", 16 * -1)
        )

        self.entry1 = tk.Text(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            font=("標楷體", 16 * -1),
            highlightthickness=0
        )

        self.entry1.place(
            x=400.0,
            y=150.0,
            width=475.0,
            height=200.0
        )

        # self.less = self.canvas.create_text(
        #     17.0,
        #     535.0,
        #     anchor="nw",
        #     text=f'LAST：',
        #     fill="#000000",
        #     font=("InriaSans Bold", 14 * -1)
        # )
        
        self.button6 = ctk.CTkButton(
            self,
            text = "確認內文",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [check_contents(),give_hint()],
            font=('Aries', 16),
            corner_radius= 30, #邊角圓弧度
            width=80.0,
            height=35.0
        )

        self.button6.place(
            x=675.0,
            y=65.0,
        )
        
        self.button5 = ctk.CTkButton(
            self,
            text = "生成內文",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [generate_contents(),give_hint(), temperature_up()],
            font=('Aries', 16),
            corner_radius= 30, #邊角圓弧度
            width=80.0,
            height=35.0,
        )

        self.button5.place(
            x=550.0,
            y=65.0,
        )
        
        self.canvas.pack()
    
class PageThree(tk.Frame):  #stable diffusion

    def __init__(self, master, *args, **kwargs):
        tk.Frame.__init__(self,master, *args, **kwargs)
        global try1, tousr, page_of_ppt, spl
        
        def finishornot():
            global spl_copy
            status = True
            if(len(spl_copy) == 1):
                status = True
            else:
                status = False
                
            if status == True:
                generate_one_page_ppt()
                finish_ppt()
                print('done')
                master.switch_Canvas(homePage)
            else:
                generate_one_page_ppt()
                master.switch_Canvas(PageTwo)

        def choose_pic0():
            global image_num, max_image, imgchoose, picc
            if picc != 0:
                imgchoose += 1
                path = r'pic0.png'
                if os.path.exists(path):
                    if os.path.exists(f'img{image_num}.png'):
                        os.remove(f'img{image_num}.png')
                    os.rename('pic0.png',f'img{image_num}.png')
                img_list.append(f'img{image_num}.png')
                self.imgw1 = PhotoImage(file='white.png')
                self.imgw1 = self.imgw1.zoom(25) #with 250, I ended up running out of memory
                self.imgw1 = self.imgw1.subsample(32)
                self.canvas.itemconfig(self.image_1,image=self.imgw1)
                image_num += 1
                max_image -= 1
                # self.canvas.itemconfig(self.less, text=f"LAST: {picc}")
                picc -= 1
                if picc == 0:
                    self.button1.configure(text = '圖片選擇已達上限')
                    self.button2.configure(text = '圖片選擇已達上限')
                    self.button4.configure(text = '圖片選擇已達上限')

        def choose_pic1():
            global image_num, max_image, imgchoose, picc
            if picc != 0:
                imgchoose += 1
                path = r'pic1.png'
                if os.path.exists(path):
                    if os.path.exists(f'img{image_num}.png'):
                        os.remove(f'img{image_num}.png')
                    os.rename('pic1.png',f'img{image_num}.png')
                img_list.append(f'img{image_num}.png')
                self.imgw2 = PhotoImage(file='white.png')
                self.imgw2 = self.imgw2.zoom(25) #with 250, I ended up running out of memory
                self.imgw2 = self.imgw2.subsample(32)
                self.canvas.itemconfig(self.image_2,image=self.imgw2)
                image_num += 1
                max_image -= 1
                # self.canvas.itemconfig(self.less, text=f"LAST: {picc}")
                picc -= 1
                if picc == 0:
                    self.button1.configure(text = '圖片選擇已達上限')
                    self.button2.configure(text = '圖片選擇已達上限')
                    self.button4.configure(text = '圖片選擇已達上限')
                    

        def choose_pic2():
            global image_num, max_image, imgchoose, picc
            if picc != 0:
                imgchoose += 1 
                path = r'pic2.png'
                if os.path.exists(path):
                    if os.path.exists(f'img{image_num}.png'):
                        os.remove(f'img{image_num}.png')
                    os.rename('pic2.png',f'img{image_num}.png')
                img_list.append(f'img{image_num}.png')
                self.imgw3 = PhotoImage(file='white.png')
                self.imgw3 = self.imgw3.zoom(25) #with 250, I ended up running out of memory
                self.imgw3 = self.imgw3.subsample(32)
                self.canvas.itemconfig(self.image_3,image=self.imgw3)
                image_num += 1
                max_image -= 1
                # self.canvas.itemconfig(self.less, text=f"LAST: {picc}")
                picc -= 1
                if picc == 0:
                    self.button1.configure(text = '圖片選擇已達上限')
                    self.button2.configure(text = '圖片選擇已達上限')
                    self.button4.configure(text = '圖片選擇已達上限')
                    
        def skip():
            global picc, img_list
            if picc == 1:
                img_list.append('white.png')
            else:
                for i in range(0, 2):
                    img_list.append('white.png')

        def generate_pict():
            global pic_loc,ansl, chatans
            chatans.clear()
            pic_loc.clear()
            ansl.clear()
            
            key_words = self.entry.get()
            prompts = f'''請幫我翻譯以下句子成英文，但如果輸入的句子皆為英文，則照原輸入輸出: 
            "{key_words}"
            '''
            print("a")
            print(key_words)
            doQA(prompts, 2048, 0.5)
            print('cha')
            print(chatans[0])
            key_words = ''.join(chatans[0])
            print("b")
            chatans.clear()
            prompts = f'''refine "{key_words}" within 30 words'''
            doQA(prompts, 2048, 0.5)
            print('cha')
            print(chatans[0])
            key_words = ''.join(chatans[0])
            print("c")
            chatans.clear()
            key_words = "please generate a ((master piece)),((high quality)), ultra-detailed illustration, related to " + key_words
 
            stability_api = client.StabilityInference(
            key=os.environ['STABILITY_KEY'], 
            verbose=True
            )
            for i in range(3):
                s = f'answer{i}' 
                r = random.randrange(0, 10000000)
                s = stability_api.generate(
                    prompt = key_words,
                    seed = r, 
                    steps = 30,
                    height= 512,
                    width=512
                )
                ansl.append(s)
                for resp in ansl[i]:
                    for artifact in resp.artifacts:
                        if artifact.finish_reason == generation.FILTER:
                            warnings.warn(
                                "Your request activated the API's safety filters and could not be processed."
                                "Please modify the prompt and try again.")
                        if artifact.type == generation.ARTIFACT_IMAGE:
                            img = Image.open(io.BytesIO(artifact.binary))
                            ps = f'pic{i}'
                            img.save(f'{ps}.png', 'png')
                            pic_loc.append(ps)

            self.img0 = PhotoImage(file='pic0.png')
            self.img0 = self.img0.zoom(25) #with 250, I ended up running out of memory
            self.img0 = self.img0.subsample(51)
            self.image_1 = self.canvas.create_image(
                142.0,
                283.0,
                image=self.img0
            )

            self.img1 = PhotoImage(file='pic1.png')
            self.img1 = self.img1.zoom(25) #with 250, I ended up running out of memory
            self.img1 = self.img1.subsample(51)
            self.image_2 = self.canvas.create_image(
                447.0,
                283.0,
                image=self.img1
            )

            self.img2 = PhotoImage(file='pic2.png')
            self.img2 = self.img2.zoom(25) #with 250, I ended up running out of memory
            self.img2 = self.img2.subsample(51)
            self.image_3 = self.canvas.create_image(
                752.0,
                283.0,
                image=self.img2
            )

            self.image_1.itemconfig(self.image_1,image=self.img0)
            self.image_2.itemconfig(self.image_2,image=self.img1)
            self.image_3.itemconfig(self.image_3,image=self.img2)

        def inittext():
            self.button1.configure(text = '選擇此圖片')
            self.button2.configure(text = '選擇此圖片')
            self.button4.configure(text = '選擇此圖片')

        def give_hint():
            global hint_yn
    
            if hint_yn == 0:
                self.canvas.delete('origin')
                
                self.canvas.create_text(
                    20.0,
                    10.0,
                    anchor="nw",
                    text="按下「選擇此圖片」選擇要放進簡報的圖片，如果都不滿意可以按「生成圖片」重新生成",
                    fill="#FFFFFF",
                    font=("InriaSans Regular", 20 * -1),
                    tags= 'hint1'
                )
                hint_yn = 1

            elif hint_yn == 1:
                self.canvas.delete('hint1')
                self.canvas.create_text(
                    20.0,
                    10.0,
                    anchor="nw",
                    text="選好後請按右上方的「完成本頁」，代表你已經完成這頁簡報囉!",
                    fill="#FFFFFF",
                    font=("InriaSans Regular", 20 * -1),
                    tags= 'hint2'
                )
                hint_yn = 1


        def hint_to_0():
            global hint_yn
            hint_yn = 0
            self.canvas.delete('hint2')
            self.canvas.create_text(
                20.0,
                10.0,
                anchor="nw",
                text="輸入想生成的圖片關鍵字後，請按下「生成圖片」按鈕",
                fill="#FFFFFF",
                font=("InriaSans Regular", 20 * -1),
                tags= 'origin'
            )
        
        self.canvas = tk.Canvas(
            self,
            bg = "#E26565",
            height = 600,
            width = 900,
            bd = 0,
            highlightthickness = 0,
            relief = "ridge"
        )
        
        self.buttonw = ctk.CTkButton(
            self,
            text = "略過",
            text_color= 'black', #文字顏色
            fg_color='#E8E9E9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [skip(), master.switch_Canvas(PageTwo), finishornot(), inittext(), hint_to_0()],
            font=('Aries', 20),
            corner_radius= 0, #邊角圓弧度
            width=80.0,
            height=35.0
        )

        self.buttonw.place(
            x=665.0,
            y=60.0
        )

        self.canvas.place(x = 0, y = 0)
        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            600.0,
            fill="#E8E9E9",
            outline=""
        )

        self.canvas.create_rectangle(
            0.0,
            0.0,
            900.0,
            50.0,
            fill="#8696A9",
            outline=""
        )

        self.canvas.create_text(
            20.0,
            10.0,
            anchor="nw",
            text="輸入想生成的圖片關鍵字後，請按下「生成圖片」按鈕",
            fill="#FFFFFF",
            font=("InriaSans Regular", 20 * -1),
            tags='origin'
        )
        
        self.img0 = PhotoImage(file='white.png')
        # self.img0 = self.img0.zoom(25) #with 250, I ended up running out of memory
        # self.img0 = self.img0.subsample(32)
        self.image_1 = self.canvas.create_image(
            20.0,
            180.0,
            image = self.img0,
            anchor='nw', 
            tags="image"
        )

        self.img1 = PhotoImage(file='white.png')
        # self.img1 = self.img1.zoom(25) #with 250, I ended up running out of memory
        # self.img1 = self.img1.subsample(32)
        self.image_2 = self.canvas.create_image(
            325.0,
            180.0,
            image = self.img1,
            anchor='nw', 
            tags="image"
        )
        
        self.img2 = PhotoImage(file='white.png')
        # self.img2 = self.img2.zoom(25) #with 250, I ended up running out of memory
        # self.img2 = self.img2.subsample(32)
        self.image_3 = self.canvas.create_image(
            630.0,
            180.0,
            image = self.img2,
            anchor='nw', 
            tags="image"
        )
        
        self.button1 = ctk.CTkButton(
            self,
            text = "選擇此圖片",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [choose_pic0(),give_hint()],
            font=('Inter Bold', 14),
            corner_radius= 0, #邊角圓弧度
            width=150.0,
            height=35.0
        )

        self.button1.place(
            x=67.0,
            y=444.0
        )

        self.button2 = ctk.CTkButton(
            self,
            text = "選擇此圖片",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [choose_pic1(),give_hint()],
            font=('Inter Bold', 14),
            corner_radius= 0, #邊角圓弧度
            width=150.0,
            height=35.0
        )

        self.button2.place(
            x=372.0,
            y=444.0
        )

        self.button3 = ctk.CTkButton(
            self,
            text = "生成圖片",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [give_hint(),generate_pict()],
            font=('Aries', 20),
            corner_radius= 30, #邊角圓弧度
            width=250.0,
            height=35.0
        )

        self.button3.place(
            x=322.0,
            y=517.0
        )

        self.button4 = ctk.CTkButton(
            self,
            text = "選擇此圖片",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [choose_pic2(),give_hint()],
            font=('Inter Bold', 14),
            corner_radius= 0, #邊角圓弧度
            width=150.0,
            height=35.0
        )

        self.button4.place(
            x=677.0,
            y=444.0
        )

        # self.less = self.canvas.create_text(
        #     17.0,
        #     524.0,
        #     anchor="nw",
        #     text=f"LAST：",
        #     fill="#000000",
        #     font=("InriaSans Bold", 14 * -1)
        # )

        self.buttonr = ctk.CTkButton(
            self,
            command=lambda: finish_ppt(),
            fg_color='red', #按鈕顏色
            text="",
            # corner_radius= 30, #邊角圓弧度
            # hover_color='#7A8797',
            bg_color="red",
            # font=('Helvetica', 25, 'bold'),
            width=20.0,
            height=10.0
        )
        
        self.buttonr.place(
            x=880.0,
            y=590.0,
        )

        self.entry = tk.Entry(
            self,
            bd=0,
            bg="#FFFFFF",
            fg="#000716",
            highlightthickness=0,
            font = ('Aries', 16)
        )

        self.entry.place(
            x=145.0,
            y=95.0,
            width=600.0,
            height=50.0
        )

        self.button6 = ctk.CTkButton(
            self,
            text = "完成本頁",
            text_color= 'white', #文字顏色
            fg_color='#8696A9', #按鈕顏色
            bg_color="#E8E9E9",
            hover_color='#7A8797',
            command=lambda: [master.switch_Canvas(PageTwo), finishornot(), inittext(), hint_to_0()],
            font=('Aries', 20),
            corner_radius= 0, #邊角圓弧度
            width=98.0,
            height=50.0
        )

        self.button6.place(
            x=802.0,
            y=0.0
        )
        
        self.canvas.create_text(
            145.0,
            60.0,
            anchor="nw",
            text="關鍵字: ",
            fill="#000000",
            font=("Inter Bold", 20 * -1),
        )
        
        self.canvas.pack()

if __name__ == "__main__":
    app = SampleApp()
    app.mainloop()



